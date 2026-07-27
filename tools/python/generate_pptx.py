"""Generate banking-parser-platform.pptx from slide content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colors ────────────────────────────────────────────────────────────────────
ORANGE      = RGBColor(0xE6, 0x51, 0x00)  # #e65100
DARK_ORANGE = RGBColor(0xBF, 0x36, 0x0C)  # #bf360c
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAY_BG     = RGBColor(0xF4, 0xF4, 0xF4)
DARK_BG     = RGBColor(0x1E, 0x1E, 0x1E)
LIGHT_TEXT  = RGBColor(0xD4, 0xD4, 0xD4)
ALT_ROW     = RGBColor(0xFE, 0xF6, 0xF2)
ROW_BORDER  = RGBColor(0xDD, 0xDD, 0xDD)
LINK_BLUE   = RGBColor(0x15, 0x65, 0xC0)

W  = Inches(13.33)   # widescreen 16:9
H  = Inches(7.5)

TOTAL = 22          # total slide count (keep in sync with main())

OUT_DIR   = os.path.join(os.path.dirname(__file__), "..", "..", "docs", "slides")
OUT_FILE  = os.path.join(OUT_DIR, "banking-parser-platform.pptx")


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def fill_slide(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height,
                text, font_size=18, bold=False,
                color: RGBColor = BLACK,
                align=PP_ALIGN.LEFT,
                italic=False,
                word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf  = txb.text_frame
    tf.word_wrap = word_wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return txb


def add_paragraph(tf, text, font_size=16, bold=False,
                  color: RGBColor = BLACK,
                  align=PP_ALIGN.LEFT,
                  italic=False,
                  space_before=Pt(4)):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = _Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = "Calibri"
    return p


def add_code_box(slide, left, top, width, height, code_text, font_size=11):
    rect = add_rect(slide, left, top, width, height, DARK_BG)
    txb  = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2))
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(font_size)
        run.font.color.rgb = LIGHT_TEXT
        run.font.name  = "Courier New"
    return txb


def add_table(slide, left, top, width, rows_data,
              col_widths=None, font_size=14):
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    row_h    = Inches(0.38)
    height   = row_h * num_rows

    table = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height).table

    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            tf   = cell.text_frame
            tf.word_wrap = True
            p    = tf.paragraphs[0]
            run  = p.add_run() if not p.runs else p.runs[0]
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            if ri == 0:                  # header row
                run.font.bold  = True
                run.font.color.rgb = WHITE
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = ORANGE
            else:
                run.font.color.rgb = BLACK
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = ALT_ROW if ri % 2 == 0 else WHITE

    return table


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)
    fill_slide(slide, ORANGE)

    # bottom white bar
    add_rect(slide, 0, H - Inches(1.4), W, Inches(1.4), RGBColor(0xFF, 0xFF, 0xFF))

    add_textbox(slide,
        Inches(0.8), Inches(1.0), W - Inches(1.6), Inches(1.4),
        "Banking Fixed-Length File Platform",
        font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(2.6), W - Inches(1.6), Inches(0.6),
        "Generating, parsing and benchmarking CODA & SWIFT MT940 files",
        font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(3.2), W - Inches(1.6), Inches(0.6),
        "across 9 Java formatter approaches via Strategy Pattern + Spring Batch",
        font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(4.2), W - Inches(1.6), Inches(0.6),
        "Wallace Espindola  ·  wallace.espindola@gmail.com",
        font_size=17, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(4.8), W - Inches(1.6), Inches(0.5),
        "linkedin.com/in/wallaceespindola   ·   github.com/wallaceespindola",
        font_size=15, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

    # page number area (bottom bar)
    add_textbox(slide,
        Inches(0.4), H - Inches(1.1), W - Inches(0.8), Inches(0.4),
        f"1 / {TOTAL}", font_size=12, color=RGBColor(0x88, 0x88, 0x88),
        align=PP_ALIGN.RIGHT)


def slide_heading(slide, title, page):
    add_rect(slide, 0, 0, W, Inches(0.78), ORANGE)
    add_textbox(slide,
        Inches(0.4), Inches(0.1), W - Inches(0.8), Inches(0.6),
        title, font_size=26, bold=True, color=WHITE)
    add_textbox(slide,
        Inches(0.4), H - Inches(0.35), W - Inches(0.8), Inches(0.3),
        f"{page} / {TOTAL}", font_size=11,
        color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.RIGHT)


def slide_problem(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Problem Statement", 2)

    add_textbox(slide,
        Inches(0.5), Inches(1.0), W - Inches(1.0), Inches(0.5),
        "Multiple Java libraries claim to support fixed-length banking file formats.",
        font_size=17, color=BLACK)
    add_textbox(slide,
        Inches(0.5), Inches(1.5), W - Inches(1.0), Inches(0.5),
        "Which one should a bank standardise on for Spring Batch workloads?",
        font_size=20, bold=True, color=ORANGE)

    add_textbox(slide,
        Inches(0.5), Inches(2.2), Inches(3), Inches(0.35),
        "Evaluation criteria:", font_size=17, bold=True, color=BLACK)

    criteria = [
        ("Correctness", "Does output conform to Febelfin / SWIFT specifications?"),
        ("Performance", "Throughput in records/second, measured with JMH"),
        ("Maintainability", "Annotation quality, layout auditability, no hidden XML"),
        ("Spring Batch fit", "Chunk-oriented reader/writer compatibility"),
        ("Supply-chain health", "Release cadence, governance, dependency weight"),
        ("Operational risk", "Support model, key-person risk, CVE history"),
    ]
    for i, (label, desc) in enumerate(criteria):
        top = Inches(2.6) + i * Inches(0.52)
        add_textbox(slide, Inches(0.5), top, Inches(2.5), Inches(0.5),
                    f"{i+1}. {label}", font_size=17, bold=True, color=DARK_ORANGE)
        add_textbox(slide, Inches(3.1), top, W - Inches(3.6), Inches(0.5),
                    f"— {desc}", font_size=16, color=BLACK)

    add_rect(slide,
        Inches(0.5), Inches(6.0), W - Inches(1.0), Inches(0.8),
        GRAY_BG)
    add_textbox(slide,
        Inches(0.7), Inches(6.1), W - Inches(1.4), Inches(0.7),
        "One codebase, 9 approaches, identical domain data, automated benchmarks.",
        font_size=16, italic=True, color=DARK_ORANGE)


def slide_architecture(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Architecture", 3)

    code = """\
Web UI (HTML/CSS/JS)
        |  HTTP REST
        v
   REST API Layer           POST /api/domain/generate
   (Spring MVC)             POST /api/batch/generate
        |                   GET  /api/benchmark/results
        v
  Spring Batch Pipeline
  DomainEntityItemReader -> FileGenerationItemProcessor -> FileOutputItemWriter
                                      |
                               StrategyResolver  (O(1) map lookup)
                     /    |    |    |    |    |    |    |    \\
                BeanIO  ff4j  VL  Bindy CamelBIO Vel  SB  SB+ff4j  SB+VL
                              |
                       18 FileGenerationStrategy implementations
                       (9 approaches x CODA + SWIFT)"""

    add_code_box(slide,
        Inches(0.4), Inches(0.95), W - Inches(0.8), Inches(5.9),
        code, font_size=13)


def slide_coda(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "CODA Format — Belgian Banking Standard", 4)

    add_textbox(slide,
        Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(0.45),
        "Febelfin CODA — fixed-width ASCII, exactly 128 characters per record",
        font_size=17, bold=True, color=DARK_ORANGE)

    rows = [
        ("Record", "Meaning"),
        ("0", "File header"),
        ("1", "Movement (debit/credit transaction)"),
        ("2", "Movement detail / free communication"),
        ("8", "Information record (closing balance)"),
        ("9", "File trailer"),
    ]
    col_w = [Inches(1.2), W - Inches(1.7)]
    add_table(slide, Inches(0.5), Inches(1.5), W - Inches(1.0),
              rows, col_w, font_size=15)

    code = """\
0310HDR       BE68539007547034                     EUR000000000000000022052622052...
1310REF0000001BE12345678901234567890123456789012345EUR000000000000075029042629042...
9000TRAILER   ..."""
    add_code_box(slide,
        Inches(0.5), Inches(4.7), W - Inches(1.0), Inches(1.25),
        code, font_size=13)

    add_textbox(slide,
        Inches(0.5), Inches(6.1), W - Inches(1.0), Inches(0.4),
        "Each field has an exact byte offset — annotations define the mapping.",
        font_size=14, italic=True, color=RGBColor(0x55, 0x55, 0x55))


def slide_swift(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "SWIFT MT940 — International Statement Format", 5)

    add_textbox(slide,
        Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(0.4),
        "SWIFT MT940 — tag-based messages, inter-message separator '---'",
        font_size=17, bold=True, color=DARK_ORANGE)

    rows = [
        ("Tag",    "Field",                  "Example"),
        (":20:",   "Transaction reference",  "STMT000001"),
        (":25:",   "Account identification", "BE68539007547034/EUR"),
        (":28C:",  "Statement / sequence",   "00001/001"),
        (":60F:",  "Opening balance",        "C260429EUR10000,00"),
        (":61:",   "Statement line",         "260429260429C750NMSCREF001"),
        (":86:",   "Narrative",              "Payment for services"),
        (":62F:",  "Closing balance",        "C260429EUR10750,00"),
    ]
    col_w = [Inches(1.2), Inches(4.0), W - Inches(5.7)]
    add_table(slide, Inches(0.5), Inches(1.5), W - Inches(1.0),
              rows, col_w, font_size=14)

    code = """\
:20:STMT000001
:25:BE68539007547034/EUR
:60F:C260429EUR10000,00
:61:260429260429C750NMSCREF001NONREF
:62F:C260429EUR10750,00
---"""
    add_code_box(slide,
        Inches(0.5), Inches(5.3), W - Inches(1.0), Inches(1.55),
        code, font_size=13)


def slide_libraries(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "9 Formatter Approaches", 6)

    rows = [
        ("Approach",         "Mechanism",                                       "CODA W", "CODA R", "SWIFT"),
        ("BeanIO",           "StreamBuilder + FieldBuilder.at() (0-based)",     "Yes",    "Yes",    "Yes"),
        ("fixedformat4j",    "@Record(length=128) + @Field(offset, length)",    "Yes",    "Yes",    "Yes"),
        ("fixedlength",      "@FixedLine + @FixedField(offset, length)",        "Yes",    "Yes",    "Yes"),
        ("Camel Bindy",      "@FixedLengthRecord + @DataField(pos, length)",    "Yes",    "Yes",    "Yes"),
        ("Camel BeanIO",     "XML stream mapping via Camel dataformat",         "Yes",    "Yes",    "Yes"),
        ("Velocity",         ".vm template files (write-only for CODA)",        "Yes",    "—",      "Yes"),
        ("Spring Batch",     "FormatterLineAggregator + FixedLengthTokenizer",  "Yes",    "Yes",    "Yes"),
        ("Spring Batch + ff4j",         "Spring Batch components, layout from @Field",      "Yes", "Yes", "Yes"),
        ("Spring Batch + fixedlength",  "Spring Batch components, layout from @FixedField", "Yes", "Yes", "Yes"),
    ]
    col_w = [Inches(2.9), Inches(5.0), Inches(1.3), Inches(1.3), Inches(1.3)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=13)

    add_textbox(slide,
        Inches(0.5), Inches(6.6), W - Inches(1.0), Inches(0.4),
        "All approaches share the same domain data and produce byte-comparable output files.",
        font_size=13, italic=True, color=RGBColor(0x55, 0x55, 0x55))


def slide_strategy(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Strategy Pattern — One Interface, 18 Implementations", 8)

    code1 = """\
public interface FileGenerationStrategy {
    String generate(List<Transaction> txs, List<Account> accounts);
    List<Transaction> parse(String fileContent);
    FileType getFileType();   // CODA | SWIFT
    Library   getLibrary();   // BEANIO | FIXFORMAT4J | FIXEDLENGTH | BINDY
                              // CAMEL_BEANIO | VELOCITY | SPRING_BATCH
                              // SPRING_BATCH_FF4J | SPRING_BATCH_FIXEDLENGTH
    default String strategyKey() { return getFileType() + "_" + getLibrary(); }
}"""
    add_code_box(slide, Inches(0.4), Inches(0.95),
                 W - Inches(0.8), Inches(2.5), code1, font_size=13)

    code2 = """\
// Resolution — O(1) map lookup, no if/switch chains
FileGenerationStrategy s = resolver.resolve(FileType.CODA, Library.SPRING_BATCH_FF4J);
String codaFile = s.generate(transactions, accounts);"""
    add_code_box(slide, Inches(0.4), Inches(3.6),
                 W - Inches(0.8), Inches(1.5), code2, font_size=13)

    add_textbox(slide,
        Inches(0.5), Inches(5.3), W - Inches(1.0), Inches(0.5),
        "StrategyResolver auto-wires all 18 beans from the Spring context at startup.",
        font_size=16, color=DARK_ORANGE)


def slide_batch(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Spring Batch Pipeline", 9)

    code = """\
bankingFileGenerationJob  (restartable — saveState=true)
+-- fileGenerationStep    (chunk-size = 100)
    +-- DomainEntityItemReader
    |     Loads all Transaction rows from H2
    +-- FileGenerationItemProcessor
    |     Resolves FileGenerationStrategy by (fileType, library)
    |     Calls strategy.generate([transaction], accounts)
    +-- FileOutputItemWriter
          Buffers chunk output; on @AfterStep writes output/*.txt
          Stores file content + metadata in step ExecutionContext
               |
               v
         BatchMetricsListener (JobExecutionListener)
         Saves BenchmarkMetrics row to H2 on job completion"""
    add_code_box(slide, Inches(0.4), Inches(0.95),
                 W - Inches(0.8), Inches(5.1), code, font_size=13)

    add_textbox(slide,
        Inches(0.5), Inches(6.25), W - Inches(1.0), Inches(0.4),
        "Job parameters: fileType, library, operationId, runTimestamp",
        font_size=14, italic=True, color=RGBColor(0x55, 0x55, 0x55))


def slide_api(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "REST API", 17)

    rows = [
        ("Method", "Endpoint",                        "Description"),
        ("POST",   "/api/domain/generate",            "Seed H2 with sample data (?loadProfile=LOW|MEDIUM|HIGH)"),
        ("POST",   "/api/batch/generate",             "Trigger Spring Batch job {fileType, library}"),
        ("GET",    "/api/batch/history",              "Last 50 job executions"),
        ("GET",    "/api/benchmark/results",          "All benchmark metrics"),
        ("GET",    "/api/benchmark/export/csv",       "Export as CSV"),
        ("GET",    "/api/benchmark/export/json",      "Export as JSON"),
        ("GET",    "/api/benchmark/export/markdown",  "Export as Markdown"),
        ("GET",    "/api/benchmark/export/html",      "Velocity-rendered HTML report"),
        ("GET",    "/actuator/health",                "Application health + version"),
        ("GET",    "/actuator/info",                  "App name, version, description"),
    ]
    col_w = [Inches(1.1), Inches(3.8), W - Inches(5.4)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=13)


def slide_benchmark(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Benchmark Metrics", 18)

    rows = [
        ("Metric",                "Description"),
        ("throughputRps",         "Records processed per second"),
        ("batchDurationMs",       "Total Spring Batch job wall-clock time"),
        ("generationDurationMs",  "File serialisation time only"),
        ("parseDurationMs",       "File parsing time (round-trip)"),
        ("symmetryRate",          "% of parsed transactions matching original domain data"),
        ("successRate",           "% of chunks completed without error"),
    ]
    col_w = [Inches(3.2), W - Inches(3.7)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=15)

    code = """\
# Run JMH benchmark suite (36 @Benchmark methods)
mvn test -Pbenchmark

# Export results
curl http://localhost:8080/api/benchmark/export/csv -o results.csv
curl http://localhost:8080/api/benchmark/export/json"""
    add_code_box(slide, Inches(0.4), Inches(4.55),
                 W - Inches(0.8), Inches(1.8), code, font_size=13)


def slide_recommendations(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Decision Guide", 16)

    rows = [
        ("Use case",                          "Pick",                        "Why"),
        ("New Spring Batch job in a bank",    "Spring Batch + fixedformat4j","Batch-native runtime, layout declared once, no extra runtime deps"),
        ("Minimal dependency footprint",      "Spring Batch + fixedlength",  "33 KB library, same annotation-driven layout"),
        ("Complex CODA grammar",              "BeanIO",                      "Richest grammar model (record groups, repeating segments)"),
        ("Camel routes already in production","Camel Bindy",                 "Native dataformat inside existing routes"),
        ("Layout auditable outside the code", "Camel BeanIO",                "XML mapping files, reviewable without Java"),
        ("Rendering statements / reports",    "Velocity",                    "Template engine, write-only by design"),
        ("No new dependency allowed",         "Spring Batch native",         "Ships with the batch runtime"),
    ]
    col_w = [Inches(3.6), Inches(3.0), W - Inches(7.4)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=12)

    add_rect(slide,
        Inches(0.4), Inches(5.5), W - Inches(0.8), Inches(1.0),
        GRAY_BG)
    add_textbox(slide,
        Inches(0.6), Inches(5.55), W - Inches(1.2), Inches(0.45),
        "Recommendation: standardise on one approach per system.",
        font_size=15, bold=True, color=DARK_ORANGE)
    add_textbox(slide,
        Inches(0.6), Inches(6.0), W - Inches(1.2), Inches(0.4),
        "Benchmark on your own hardware, then pin the version and treat the layout model as a controlled artefact.",
        font_size=13, italic=True, color=BLACK)


def slide_quality(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Code Quality & CI/CD", 19)

    add_textbox(slide,
        Inches(0.5), Inches(0.95), Inches(4), Inches(0.4),
        "Testing — 149 tests across 13 test classes:",
        font_size=17, bold=True, color=BLACK)

    rows = [
        ("Category",   "Tests",                                              "Coverage"),
        ("Unit",       "DomainDataGeneratorTest, CodaRecordTest, AnnotatedLayoutTest", "Mock repos, fields, layout reflection"),
        ("Integration","StrategyResolverTest, CodaStrategyTest, SwiftStrategyTest", "All 18 strategies"),
        ("Symmetry",   "SymmetryTest",                                        "Round-trip: generate -> parse -> compare"),
        ("Golden file","GoldenFileTest",                                      "128-char CODA lines, MT940 tags"),
        ("API",        "DomainControllerTest, BatchControllerTest",           "MockMvc"),
        ("Actuator",   "ActuatorTest, SwaggerAvailabilityTest",               "TestRestTemplate"),
    ]
    col_w = [Inches(1.5), Inches(4.5), W - Inches(6.5)]
    add_table(slide, Inches(0.4), Inches(1.45), W - Inches(0.8),
              rows, col_w, font_size=13)

    bullets = [
        "CI/CD: GitHub Actions — build · test · benchmark · CodeQL · release",
        "Coverage: JaCoCo enforced at minimum threshold · Dependabot weekly PRs",
    ]
    for i, b in enumerate(bullets):
        add_textbox(slide,
            Inches(0.5), Inches(5.55) + i * Inches(0.5),
            W - Inches(1.0), Inches(0.45),
            f"• {b}", font_size=15, color=BLACK)


def slide_quickstart(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Quick Start", 20)

    code = """\
# Clone and build (Java 21 + Maven 3.9 required — no Node.js needed)
git clone https://github.com/wallaceespindola/fixed-length-converters
cd fixed-length-converters
mvn clean install

# Start in dev mode (Swagger UI at /swagger-ui.html)
mvn spring-boot:run -Dspring-boot.run.profiles=dev

# Step 1 -- Generate domain data
curl -X POST http://localhost:8080/api/domain/generate?loadProfile=HIGH

# Step 2 -- Run batch job (pick any approach)
curl -X POST http://localhost:8080/api/batch/generate \\
  -H "Content-Type: application/json" \\
  -d '{"fileType":"CODA","library":"SPRING_BATCH_FF4J"}'

# Step 3 -- Export benchmark results
curl http://localhost:8080/api/benchmark/export/csv -o results.csv"""

    add_code_box(slide, Inches(0.4), Inches(0.95),
                 W - Inches(0.8), Inches(5.95), code, font_size=13)


def slide_stack(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Technology Stack", 21)

    rows = [
        ("Area",        "Technology"),
        ("Language",    "Java 21"),
        ("Backend",     "Spring Boot 3.4.5, Spring Batch 5.2.2, Spring Data JPA"),
        ("Database",    "H2 In-Memory"),
        ("API Docs",    "OpenAPI + Swagger UI (dev profile)"),
        ("Monitoring",  "Spring Actuator (/health, /info, version indicator)"),
        ("Frontend",    "Vanilla HTML/CSS/JS + Chart.js"),
        ("Build",       "Maven (single mvn clean install, no profiles)"),
        ("Testing",     "JUnit 5 + Mockito, 149 tests, 36 JMH benchmarks"),
        ("CI/CD",       "GitHub Actions (build, test, benchmark, CodeQL)"),
    ]
    col_w = [Inches(2.5), W - Inches(3.0)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=15)


def slide_thankyou(prs):
    slide = blank_slide(prs)
    fill_slide(slide, ORANGE)

    add_rect(slide, 0, H - Inches(1.4), W, Inches(1.4), WHITE)

    add_textbox(slide,
        Inches(0.8), Inches(1.2), W - Inches(1.6), Inches(1.2),
        "Thank You",
        font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(2.7), W - Inches(1.6), Inches(0.7),
        "github.com/wallaceespindola/fixed-length-converters",
        font_size=22, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(3.6), W - Inches(1.6), Inches(0.55),
        "Wallace Espindola",
        font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(4.15), W - Inches(1.6), Inches(0.5),
        "wallace.espindola@gmail.com",
        font_size=17, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(4.7), W - Inches(1.6), Inches(0.5),
        "linkedin.com/in/wallaceespindola   ·   github.com/wallaceespindola",
        font_size=15, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.8), Inches(5.5), W - Inches(1.6), Inches(0.5),
        "Questions welcome — slides, code and benchmarks all open source",
        font_size=15, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide,
        Inches(0.4), H - Inches(1.1), W - Inches(0.8), Inches(0.4),
        f"{TOTAL} / {TOTAL}", font_size=12,
        color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.RIGHT)



def slide_annotated_layout(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Annotation-Derived Layouts — Fixing the Slicing Problem", 7)

    add_textbox(slide,
        Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(0.4),
        "Native Spring Batch needs the layout twice: a Range list to read, a format string to write.",
        font_size=16, color=BLACK)
    add_textbox(slide,
        Inches(0.5), Inches(1.35), W - Inches(1.0), Inches(0.4),
        "Two copies of the same offsets drift apart silently.",
        font_size=16, bold=True, color=ORANGE)

    code = """\
// Before -- layout restated by hand, 12 magic ranges
t.setColumns(new Range(1,1), new Range(2,4), new Range(5,14), /* ... */);

// After -- layout read from the annotated model, once
AnnotatedLayout layout = AnnotatedLayout.fromFixedFormat4j(Ff4jCodaRecord.class);
FixedLengthTokenizer tokenizer = layout.tokenizer();   // read path
String format = layout.formatString();                 // write path"""
    add_code_box(slide, Inches(0.4), Inches(1.9),
                 W - Inches(0.8), Inches(2.5), code, font_size=13)

    bullets = [
        "Offsets, widths, alignment and padding char come from @Field / @FixedField",
        "Constructor validates gaps, overlaps and the 128-char total — layout errors fail fast",
        "Output is byte-identical to hand-sliced Spring Batch (asserted in AnnotatedLayoutTest)",
    ]
    for i, b in enumerate(bullets):
        add_textbox(slide, Inches(0.5), Inches(4.6) + i * Inches(0.5),
                    W - Inches(1.0), Inches(0.45),
                    f"• {b}", font_size=15, color=BLACK)


def slide_health(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Library Health — Verified 2026-07-27", 10)

    rows = [
        ("Library",        "Coordinates (groupId:artifactId)",                    "Pinned", "Latest", "Released",   "Repo activity"),
        ("BeanIO",         "com.github.beanio:beanio",                            "3.2.1",  "3.2.1",  "2025-02-07", "2025-02-07"),
        ("fixedformat4j",  "com.ancientprogramming.fixedformat4j:fixedformat4j",  "1.9.1",  "1.9.1",  "2026-06-17", "2026-07-25"),
        ("fixedlength",    "name.velikodniy.vitaliy:fixedlength",                 "0.15",   "0.15",   "2026-02-26", "2026-02-26"),
        ("Camel Bindy",    "org.apache.camel:camel-bindy",                        "4.21.0", "4.21.0", "2026-06-27", "daily"),
        ("Camel BeanIO",   "org.apache.camel:camel-beanio",                       "4.21.0", "4.21.0", "2026-06-27", "daily"),
        ("Velocity",       "org.apache.velocity:velocity-engine-core",            "2.4.1",  "2.4.1",  "2024-10-14", "2026-06-14"),
        ("Spring Batch",   "org.springframework.batch:spring-batch-core",         "5.2.2",  "6.0.4",  "2026-06-10", "2026-07-23"),
    ]
    col_w = [Inches(1.9), Inches(4.6), Inches(1.0), Inches(1.0), Inches(1.5), Inches(1.6)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=11)

    notes = [
        "Sources: Maven Central maven-metadata.xml + artifact Last-Modified header, GitHub REST API",
        "Spring Batch 5.2.2 is what Spring Boot 3.4.5 resolves; 6.x requires Boot 4.x",
    ]
    for i, n in enumerate(notes):
        add_textbox(slide, Inches(0.5), Inches(4.5) + i * Inches(0.42),
                    W - Inches(1.0), Inches(0.4),
                    f"• {n}", font_size=13, italic=True,
                    color=RGBColor(0x55, 0x55, 0x55))


def slide_adoption(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Adoption & Governance", 11)

    rows = [
        ("Library",       "GitHub repo",                 "Stars", "Governance",                          "Dependents", "License"),
        ("BeanIO",        "beanio/beanio",               "68",    "Community fork of the 2014 org.beanio","10",        "Apache-2.0"),
        ("fixedformat4j", "jeyben/fixedformat4j",        "52",    "Single maintainer, active",            "2",         "Apache-2.0"),
        ("fixedlength",   "g0ddest/fixedlength",         "23",    "Single maintainer, 0.x versioning",    "0",         "Apache-2.0"),
        ("Camel Bindy",   "apache/camel",                "6 273", "Apache Software Foundation",           "4",         "Apache-2.0"),
        ("Camel BeanIO",  "apache/camel",                "6 273", "Apache Software Foundation",           "4",         "Apache-2.0"),
        ("Velocity",      "apache/velocity-engine",      "413",   "Apache Software Foundation",           "1 560",     "Apache-2.0"),
        ("Spring Batch",  "spring-projects/spring-batch","2 947", "Broadcom/VMware, commercial support",  "40",        "Apache-2.0"),
    ]
    col_w = [Inches(1.8), Inches(2.9), Inches(0.9), Inches(4.0), Inches(1.4), Inches(1.5)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=11)

    add_rect(slide, Inches(0.4), Inches(4.6), W - Inches(0.8), Inches(1.15), GRAY_BG)
    add_textbox(slide, Inches(0.6), Inches(4.65), W - Inches(1.2), Inches(0.45),
        "Dependents = deps.dev dependent packages for the pinned version only — a directional proxy.",
        font_size=13, color=BLACK)
    add_textbox(slide, Inches(0.6), Inches(5.1), W - Inches(1.2), Inches(0.55),
        "Maven Central publishes no public download counts; treat any 'downloads' figure elsewhere as an estimate.",
        font_size=13, bold=True, color=DARK_ORANGE)


def slide_supply_chain(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Supply-Chain Weight", 12)

    rows = [
        ("Approach",           "Own jar", "Transitive cost",                      "Notes"),
        ("BeanIO",             "430 KB",  "none",                                 "Self-contained"),
        ("fixedformat4j",      "125 KB",  "none",                                 "Smallest annotation-driven option"),
        ("fixedlength",        "33 KB",   "none",                                 "Tiny; 0.x API stability caveat"),
        ("Camel Bindy",        "171 KB",  "camel-support, camel-api, icu4j 14 MB","45 Camel artifacts on the tree"),
        ("Camel BeanIO",       "27 KB",   "Camel core + BeanIO 2.x",              "Two ecosystems in one path"),
        ("Velocity",           "503 KB",  "commons-lang3, slf4j",                 "Templates are executable code"),
        ("Spring Batch (x3)",  "0 KB",    "already on the classpath",             "Batch runtime is a given"),
    ]
    col_w = [Inches(2.6), Inches(1.3), Inches(4.4), W - Inches(8.8)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=12)

    add_rect(slide, Inches(0.4), Inches(4.6), W - Inches(0.8), Inches(1.2), GRAY_BG)
    add_textbox(slide, Inches(0.6), Inches(4.68), W - Inches(1.2), Inches(0.45),
        "Security: Velocity < 2.3 carried CVE-2020-13936 (template -> RCE). Pinned 2.4.1 is not affected.",
        font_size=13, bold=True, color=DARK_ORANGE)
    add_textbox(slide, Inches(0.6), Inches(5.15), W - Inches(1.2), Inches(0.5),
        "The risk class remains: .vm templates must be version-controlled and never user-supplied.",
        font_size=13, italic=True, color=BLACK)


def slide_suitability(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Bank Suitability Matrix", 13)

    rows = [
        ("Approach",            "Grammar",  "Layout audit",   "Batch fit", "Support model",         "Bank verdict"),
        ("BeanIO",              "High",     "Builder code",   "Good",      "Community only",        "Good for complex CODA; key-person risk"),
        ("fixedformat4j",       "Low",      "Annotations",    "Excellent", "Single maintainer",     "Strong for simple fixed layouts"),
        ("fixedlength",         "Low",      "Annotations",    "Good",      "Single maintainer, 0.x","Prototyping, not core payments"),
        ("Camel Bindy",         "Medium",   "Annotations",    "Medium",    "ASF",                   "Only if Camel already in production"),
        ("Camel BeanIO",        "High",     "XML mappings",   "Medium",    "ASF",                   "Auditable XML, heavy dependency path"),
        ("Velocity",            "N/A",      "Templates",      "Low",       "ASF",                   "Report rendering, never parsing"),
        ("Spring Batch native", "Medium",   "Code (Range)",   "Native",    "Broadcom/VMware",       "Safe default; layout duplicated by hand"),
        ("Spring Batch + ff4j", "Medium",   "Annotations",    "Native",    "Broadcom + maintainer", "Best overall fit for a bank"),
        ("Spring Batch + f.l.", "Medium",   "Annotations",    "Native",    "Broadcom + maintainer", "Same shape, smaller dependency"),
    ]
    col_w = [Inches(2.4), Inches(1.0), Inches(1.5), Inches(1.0), Inches(2.3), W - Inches(9.0)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=11)


def slide_perf_jmh(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Measured Performance — JMH", 14)

    rows = [
        ("Approach",                    "CODA write", "CODA read", "SWIFT write", "SWIFT read"),
        ("fixedformat4j",               "8 261",      "14 198",    "11 418",      "14 433"),
        ("fixedlength",                 "6 374",      "7 810",     "11 616",      "13 819"),
        ("Velocity",                    "5 473",      "8 947",     "4 110",       "14 121"),
        ("BeanIO",                      "5 228",      "2 818",     "11 373",      "14 499"),
        ("Camel BeanIO",                "4 942",      "2 889",     "11 416",      "14 267"),
        ("Spring Batch native",         "4 642",      "9 188",     "11 615",      "13 009"),
        ("Spring Batch + ff4j",         "4 320",      "8 656",     "11 510",      "14 217"),
        ("Spring Batch + fixedlength",  "4 218",      "8 956",     "11 497",      "14 117"),
        ("Camel Bindy",                 "3 186",      "2 179",     "11 384",      "14 291"),
    ]
    col_w = [Inches(3.6), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    add_table(slide, Inches(0.4), Inches(0.95), W - Inches(0.8),
              rows, col_w, font_size=12)

    bullets = [
        "fixedformat4j leads CODA — ~1.3x write and ~1.5x read versus the next best",
        "Annotation-derived layout is free at runtime — hybrids sit within noise of native Spring Batch",
        "SWIFT converges (~11k write / ~13-14k read): shared SwiftMtRecord path; only Velocity stands out",
        "ops/s, median of 3 runs, Java 21 on Apple Silicon — gaps under ~20% are noise; measure on a quiet machine",
    ]
    for i, b in enumerate(bullets):
        add_textbox(slide, Inches(0.5), Inches(5.15) + i * Inches(0.42),
                    W - Inches(1.0), Inches(0.4),
                    f"• {b}", font_size=13, color=BLACK)


def slide_perf_pipeline(prs):
    slide = blank_slide(prs)
    slide_heading(slide, "Measured Performance — Batch Pipeline", 15)

    add_textbox(slide,
        Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(0.45),
        "End-to-end job, MEDIUM profile (100 accounts / 1 000 transactions), median of 4 warm runs — records/second",
        font_size=15, bold=True, color=DARK_ORANGE)

    rows = [
        ("Approach",                   "CODA",     "SWIFT"),
        ("BeanIO",                     "~72 000",  "~134 000"),
        ("Spring Batch + fixedlength", "~71 000",  "~167 000"),
        ("fixedformat4j",              "~68 000",  "~146 000"),
        ("Spring Batch native",        "~63 000",  "~71 000"),
        ("Camel BeanIO",               "~63 000",  "~143 000"),
        ("Spring Batch + ff4j",        "~59 000",  "~167 000"),
        ("fixedlength",                "~56 000",  "~143 000"),
        ("Camel Bindy",                "~29 000",  "~83 000"),
        ("Velocity",                   "~18 000",  "~19 000"),
    ]
    col_w = [Inches(4.5), Inches(2.5), Inches(2.5)]
    add_table(slide, Inches(0.4), Inches(1.5), W - Inches(0.8),
              rows, col_w, font_size=12)

    add_textbox(slide,
        Inches(0.5), Inches(5.6), W - Inches(1.0), Inches(0.45),
        "Every approach clears 1 000 records in under 60 ms end-to-end; most land between 14 ms and 18 ms,",
        font_size=13, italic=True, color=RGBColor(0x55, 0x55, 0x55))
    add_textbox(slide,
        Inches(0.5), Inches(6.0), W - Inches(1.0), Inches(0.45),
        "where millisecond timer resolution dominates — use the JMH numbers to compare formatters.",
        font_size=13, italic=True, color=RGBColor(0x55, 0x55, 0x55))


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = new_prs()

    slide_title(prs)                 # 1
    slide_problem(prs)               # 2
    slide_architecture(prs)          # 3
    slide_coda(prs)                  # 4
    slide_swift(prs)                 # 5
    slide_libraries(prs)             # 6
    slide_annotated_layout(prs)      # 7
    slide_strategy(prs)              # 8
    slide_batch(prs)                 # 9
    slide_health(prs)                # 10
    slide_adoption(prs)              # 11
    slide_supply_chain(prs)          # 12
    slide_suitability(prs)           # 13
    slide_perf_jmh(prs)              # 14
    slide_perf_pipeline(prs)         # 15
    slide_recommendations(prs)       # 16 — decision guide
    slide_api(prs)                   # 17
    slide_benchmark(prs)             # 18
    slide_quality(prs)               # 19
    slide_quickstart(prs)            # 20
    slide_stack(prs)                 # 21
    slide_thankyou(prs)              # 22

    prs.save(OUT_FILE)
    print(f"Saved: {OUT_FILE}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
